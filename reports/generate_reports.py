#!/usr/bin/env python3
"""Generate NemoClaw report as PDF and PowerPoint presentation."""

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))


def create_pdf():
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Title page
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 24)
    pdf.cell(0, 40, "", ln=True)
    pdf.cell(0, 15, "NVIDIA NemoClaw", ln=True, align="C")
    pdf.set_font("Helvetica", "", 14)
    pdf.cell(0, 10, "Technical Report for Management & Security Review", ln=True, align="C")
    pdf.cell(0, 10, "AI Agent Secure Execution Platform", ln=True, align="C")
    pdf.cell(0, 20, "", ln=True)
    pdf.set_font("Helvetica", "", 11)
    pdf.cell(0, 8, "License: Apache License 2.0 (Free, Open Source)", ln=True, align="C")
    pdf.cell(0, 8, "Status: Alpha (March 2026)", ln=True, align="C")
    pdf.cell(0, 8, "Source: github.com/NVIDIA/NemoClaw", ln=True, align="C")

    # Section 1: Overview
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "1. Overview", ln=True)
    pdf.set_font("Helvetica", "", 11)
    pdf.multi_cell(0, 6,
        "NVIDIA NemoClaw is an open-source reference stack (Apache 2.0) that enables "
        "secure execution of AI agents (OpenClaw) inside NVIDIA OpenShell - an isolated "
        "runtime environment with multi-layer protection.\n\n"
        "Key concept: NemoClaw wraps OpenClaw with an enterprise security layer providing "
        "sandbox isolation (Landlock/seccomp/netns), network traffic control, inference "
        "routing through a gateway, and full agent action auditing."
    )

    # Section 2: Technology Stack
    pdf.ln(8)
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "2. Technology Stack", ln=True)
    pdf.set_font("Helvetica", "", 10)

    stack_items = [
        ("NemoClaw Plugin", "TypeScript, Commander.js", "CLI interface, OpenClaw command registration"),
        ("NemoClaw Blueprint", "Python", "OpenShell resource orchestration"),
        ("OpenShell", "Go/Rust (NVIDIA)", "Sandbox: Landlock LSM, seccomp, network namespaces"),
        ("OpenClaw", "TypeScript/Node.js", "AI agent framework"),
        ("Container Image", "Docker", "Pre-built sandbox with OpenClaw installed"),
        ("Inference", "NVIDIA Cloud/NIM/vLLM/Ollama", "Routed through OpenShell gateway"),
    ]
    for name, tech, purpose in stack_items:
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(50, 7, name)
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(55, 7, tech)
        pdf.cell(0, 7, purpose, ln=True)

    # Section 3: Security
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "3. Security Architecture (4 Layers)", ln=True)
    pdf.set_font("Helvetica", "", 11)

    security_layers = [
        ("Layer 1: Docker Container", "Process and filesystem isolation"),
        ("Layer 2: Landlock LSM", "File access restrictions (read-only /usr, /lib, /etc)"),
        ("Layer 3: seccomp", "System call filtering"),
        ("Layer 4: Network Namespaces", "Network isolation with strict-by-default policy"),
    ]
    for layer, desc in security_layers:
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(0, 8, layer, ln=True)
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(0, 7, f"  {desc}", ln=True)

    pdf.ln(5)
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "Network Control", ln=True)
    pdf.set_font("Helvetica", "", 10)
    pdf.multi_cell(0, 6,
        "- Strict-by-default: all traffic blocked except explicitly allowed endpoints\n"
        "- Baseline policy in openclaw-sandbox.yaml\n"
        "- Pre-approved: NVIDIA API, GitHub, npm registry\n"
        "- Inference only through OpenShell gateway (no direct egress)\n"
        "- Operator Approval Flow: real-time TUI for approving/denying requests"
    )

    pdf.ln(5)
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "Filesystem Restrictions", ln=True)
    pdf.set_font("Helvetica", "", 10)
    fs_items = [
        ("/sandbox, /tmp, /dev/null", "Read-Write"),
        ("/usr, /lib, /proc, /etc, /app, /var/log", "Read-Only"),
        ("All other paths", "Blocked"),
    ]
    for path, access in fs_items:
        pdf.cell(100, 7, path)
        pdf.cell(0, 7, access, ln=True)

    # Section 4: Inference
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "4. Inference Routing", ln=True)
    pdf.set_font("Helvetica", "", 11)
    pdf.multi_cell(0, 6,
        "Inference requests from the agent never leave the sandbox directly. "
        "OpenShell intercepts them and routes to the configured provider:\n\n"
        "Agent (sandbox) -> OpenShell Gateway -> Provider\n\n"
        "Supported providers:\n"
        "- NVIDIA Cloud API (Nemotron models) - default\n"
        "- Local NIM Service (NVIDIA Inference Microservice)\n"
        "- vLLM on localhost (recommended for production)\n"
        "- Ollama (recommended for development/testing)"
    )

    # Section 5: Requirements
    pdf.ln(8)
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "5. System Requirements", ln=True)
    pdf.set_font("Helvetica", "", 10)
    reqs = [
        ("CPU", "4 vCPU (minimum)", "4+ vCPU (recommended)"),
        ("RAM", "8 GB (minimum)", "16 GB (recommended)"),
        ("Disk", "20 GB free (minimum)", "40 GB free (recommended)"),
        ("OS", "Ubuntu 22.04 LTS+", ""),
        ("Node.js", "20+", ""),
        ("Docker", "Installed & running", ""),
        ("OpenShell", "Installed", ""),
    ]
    for name, val1, val2 in reqs:
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(30, 7, name)
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(60, 7, val1)
        pdf.cell(0, 7, val2, ln=True)

    # Section 6: Air-gapped deployment
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "6. Air-Gapped Deployment Components", ln=True)
    pdf.set_font("Helvetica", "", 11)

    components = [
        ("NemoClaw", "github.com/NVIDIA/NemoClaw", "TypeScript plugin + Python blueprint"),
        ("OpenShell", "github.com/NVIDIA/OpenShell", "Sandbox runtime"),
        ("OpenClaw", "github.com/openclaw/openclaw", "AI agent framework"),
        ("Docker", "docker.com", "Container runtime"),
        ("Node.js 20+", "nodejs.org", "JS runtime"),
        ("Python 3.10+", "python.org", "Blueprint runtime"),
        ("Sandbox Image", "ghcr.io/nvidia/openshell-community", "Pre-built container"),
        ("vLLM or Ollama", "github.com/vllm-project/vllm", "Local inference server"),
        ("AI Model", "e.g. Nemotron, Qwen3, MiniMax M2.7", "LLM weights"),
    ]
    for name, source, note in components:
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(40, 7, name)
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(75, 7, source)
        pdf.cell(0, 7, note, ln=True)

    # Section 7: Recommendations
    pdf.ln(8)
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "7. Recommendations", ln=True)
    pdf.set_font("Helvetica", "", 11)
    pdf.multi_cell(0, 6,
        "1. PILOT: Deploy NemoClaw + Ollama + Qwen3-Coder on a test server\n"
        "2. PRODUCTION: NemoClaw + vLLM + Nemotron/Qwen3-Coder on GPU server\n"
        "3. AIR-GAPPED: Prepare all Docker images, npm/pip packages, and model weights; "
        "configure local inference routing\n"
        "4. MONITORING: Use built-in OpenShell audit + SIEM integration\n"
        "5. ALTERNATIVE: Consider OpenHands for model-agnostic deployment if NVIDIA "
        "ecosystem is not preferred"
    )

    out_path = os.path.join(OUTPUT_DIR, "NemoClaw_Report.pdf")
    pdf.output(out_path)
    print(f"PDF saved: {out_path}")


def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x76, 0xB9, 0x00)
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p2.alignment = PP_ALIGN.CENTER
        return slide

    def add_content_slide(title, bullets, note=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x76, 0xB9, 0x00)

        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        tf2 = body.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.space_after = Pt(8)

        if note:
            p3 = tf2.add_paragraph()
            p3.text = note
            p3.font.size = Pt(14)
            p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            p3.font.italic = True

        return slide

    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x2F)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x76, 0xB9, 0x00)

        cols = len(headers)
        total_rows = len(rows) + 1
        table_width = Inches(11.5)
        col_width = table_width // cols

        table_shape = slide.shapes.add_table(
            total_rows, cols,
            Inches(0.8), Inches(1.5),
            table_width, Inches(0.5 * total_rows)
        )
        table = table_shape.table

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x76, 0xB9, 0x00)

        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45) if ri % 2 == 0 else RGBColor(0x1B, 0x1B, 0x2F)

        return slide

    # Slide 1: Title
    add_title_slide(
        "NVIDIA NemoClaw",
        "Secure AI Agent Execution Platform\nTechnical Review for Management & Security"
    )

    # Slide 2: What is NemoClaw
    add_content_slide("What is NemoClaw?", [
        "Open-source reference stack by NVIDIA (Apache 2.0 license)",
        "Secure wrapper around OpenClaw AI agent framework",
        "4-layer sandbox isolation (Docker + Landlock + seccomp + netns)",
        "Strict-by-default network policy with operator approval",
        "Supports local inference: vLLM, Ollama, NVIDIA NIM",
        "Current status: Alpha (March 2026)",
        "GitHub: 14,300+ stars",
    ])

    # Slide 3: Architecture
    add_content_slide("Architecture", [
        "NemoClaw CLI (TypeScript plugin) -> registers commands in OpenClaw",
        "NemoClaw Blueprint (Python) -> orchestrates OpenShell resources",
        "NVIDIA OpenShell -> provides sandbox runtime environment",
        "Sandbox Container -> runs OpenClaw agent in isolation",
        "OpenShell Gateway -> routes inference, intercepts network requests",
        "Inference Provider -> NVIDIA Cloud / NIM / vLLM / Ollama",
        "",
        "Agent -> Sandbox -> Gateway -> Inference Provider",
    ])

    # Slide 4: Security - 4 Layers
    add_table_slide(
        "Security: 4 Layers of Isolation",
        ["Layer", "Technology", "Protection"],
        [
            ["1. Container", "Docker", "Process & filesystem isolation"],
            ["2. File Access", "Landlock LSM", "Read-only /usr, /lib, /etc; RW only /sandbox, /tmp"],
            ["3. Syscalls", "seccomp", "System call filtering"],
            ["4. Network", "netns + iptables", "Strict-by-default; operator approval for new endpoints"],
        ]
    )

    # Slide 5: Network Security
    add_content_slide("Network Security Details", [
        "All outbound traffic BLOCKED by default",
        "Only explicitly whitelisted endpoints are allowed",
        "Pre-approved: NVIDIA API, GitHub, npm registry",
        "Inference requests route through OpenShell gateway only",
        "Agent cannot directly reach external inference providers",
        "",
        "Operator Approval Flow:",
        "  1. Agent requests unknown endpoint -> blocked",
        "  2. Request displayed in TUI for operator review",
        "  3. Operator approves/denies in real-time",
        "  4. Approved endpoint added to session policy",
    ])

    # Slide 6: Inference Routing
    add_table_slide(
        "Supported Inference Providers",
        ["Provider", "Description", "Use Case"],
        [
            ["NVIDIA Cloud API", "Nemotron models (default)", "Cloud deployment"],
            ["NVIDIA NIM", "Local inference microservice", "Enterprise GPU servers"],
            ["vLLM", "Open-source inference server", "Production (recommended)"],
            ["Ollama", "Lightweight local model server", "Development/testing"],
        ]
    )

    # Slide 7: System Requirements
    add_table_slide(
        "System Requirements",
        ["Resource", "Minimum", "Recommended"],
        [
            ["CPU", "4 vCPU", "4+ vCPU"],
            ["RAM", "8 GB", "16 GB"],
            ["Disk", "20 GB free", "40 GB free"],
            ["OS", "Ubuntu 22.04+", "Ubuntu 22.04+"],
            ["Node.js", "20+", "20+"],
            ["Docker", "Installed", "Installed"],
            ["OpenShell", "Installed", "Installed"],
        ]
    )

    # Slide 8: Air-gapped Components
    add_content_slide("Components for Air-Gapped Deployment", [
        "Required Software:",
        "  - NemoClaw (GitHub: NVIDIA/NemoClaw)",
        "  - OpenShell (GitHub: NVIDIA/OpenShell)",
        "  - OpenClaw (GitHub: openclaw/openclaw)",
        "  - Docker, Node.js 20+, Python 3.10+",
        "  - Sandbox container image (ghcr.io/nvidia/openshell-community)",
        "",
        "Inference (choose one):",
        "  - vLLM (production) or Ollama (dev/test)",
        "",
        "Models: Nemotron-3-Super-120B, Qwen3-Coder-30B, MiniMax M2.7",
    ])

    # Slide 9: Comparison
    add_table_slide(
        "Comparison with Alternatives",
        ["Solution", "Local Models", "Sandbox", "License", "Maturity"],
        [
            ["NemoClaw", "Yes (vLLM/Ollama)", "4-layer", "Apache 2.0", "Alpha"],
            ["OpenHands", "Yes (Ollama/vLLM)", "Docker/K8s", "MIT", "Mature (65K stars)"],
            ["OpenCode", "Yes (Ollama)", "None", "Open Source", "Mature"],
            ["Aider", "Yes (Ollama)", "None", "Apache 2.0", "Mature (40K stars)"],
            ["NanoClaw", "As tool only", "Docker", "MIT", "Medium"],
        ]
    )

    # Slide 10: Recommendations
    add_content_slide("Recommendations", [
        "1. PILOT: Deploy NemoClaw + Ollama + Qwen3-Coder on test server",
        "2. PRODUCTION: NemoClaw + vLLM + Nemotron on GPU server",
        "3. AIR-GAPPED: Pre-package all images, packages, and model weights",
        "4. MONITORING: Use OpenShell audit + SIEM integration",
        "",
        "Risks to consider:",
        "  - Alpha status: APIs may change (pin versions)",
        "  - Local model quality: test for specific use cases",
        "  - Deployment complexity: use Ansible/Terraform automation",
    ])

    # Slide 11: Next Steps
    add_content_slide("Next Steps", [
        "1. Approve pilot project scope and timeline",
        "2. Allocate test server (4+ vCPU, 16GB RAM, GPU optional)",
        "3. Deploy NemoClaw + Ollama with Qwen3-Coder model",
        "4. Security team: review sandbox policies and network rules",
        "5. Run controlled tests with sample coding tasks",
        "6. Evaluate results and decide on production deployment",
        "",
        "Contact: github.com/NVIDIA/NemoClaw",
        "Documentation: docs.nvidia.com/nemoclaw/latest",
    ])

    out_path = os.path.join(OUTPUT_DIR, "NemoClaw_Presentation.pptx")
    prs.save(out_path)
    print(f"PPTX saved: {out_path}")


if __name__ == "__main__":
    create_pdf()
    create_pptx()
